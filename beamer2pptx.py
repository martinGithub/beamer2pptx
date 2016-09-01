
# GOAL
# 	the goal of this script is to help automatize the process of converting 
#	a Beamer Presentation (Latex) into a microsoft powerpoint
# 	it parses the file given in the variables texfile using regular
#	expressions and creates a powerpoint document with the
# 	same number of slides, with the slides titles , the images in each slides (but does not keep the size or layaout yet)
# 	and creates an image for each equation that are also added to the slides
#
# USAGE
# 	replace the right part in line texfile='presentation_test.tex' using the name of you latex file
#
# LIMITATIONS
# 	parsing the latex robustly with regular expression seems difficult in case of nested brackets etc
# 	we should use a actual parser or maybe convert the document to xml using LaTexXML ?  http://dlmf.nist.gov/LaTeXML/get.html 
# 	could first convert to html using the commande htlatex and then parse the html# 
# 	but on my machine i get errors when running htlatex (Can't find/open #	file `mathkerncmssi8.tfm')
# 	i get the same error using tex4ht 
#	i did not get the time to pin down the problem 
# 	plasTeX fails parsing when the latex uses 
# 		-\usepackage{mathtools}
#  		-french accents 
#

texfile='example.tex'
outputfile='example.pptx'


from pptx import Presentation
import re
import shutil
from pptx.util import Inches
import os, requests


def formula_as_file( formula, file, negate=False,header='' ):
    
    latexfile = open('tmp_equation.tex', 'w')
    latexfile.write('\\documentclass[preview]{standalone}')    
    #latexfile.write('\\input{header.tex}') # uncomment this line and change false to true line 87 if you have macros that need to be use in the original latex document
    latexfile.write('\n\\begin{document}')   
    latexfile.write('$%s$'%formula)
    latexfile.write('\n\\end{document}  ') 
    latexfile.close()
    os.system( 'pdflatex tmp_equation.tex %s'%file )
    os.system( 'convert -density 600  tmp_equation.pdf -quality 90  %s'%file )
   


    
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]


title.text = "Hello"
subtitle.text = "beamer to pptx conversion test"

bullet_slide_layout = prs.slide_layouts[1]

raw = open(texfile)
filecontent=raw.read()



# Experiment using plasTeX
# 	Instantiate a TeX processor and parse the input text
#try:
#    from plasTeX.TeX import TeX
#    tex = TeX()
#    tex.input(filecontent)
#    document = tex.parse()
#except:
#    print 'could not parse the latex fil eusing plasTex'
#    raise
    


filecontent=re.sub(r"""^%.*\n""",'%', filecontent)# remove comments
filecontent=re.sub(r"""%.*""",'%', filecontent)# remove comments
m=re.findall(r"""\\documentclass\{.*?\}(.*?)\\begin\{document\}""",filecontent ,re.DOTALL)
header=  m[0]

if False:
    header=re.sub(re.compile(r"""\\usetheme\[.*?\]\{.*?\}""",re.DOTALL),r'', header,)
    header=re.sub(re.compile(r"""\\title\[.*?\]\{.*?\}""",re.DOTALL),r'', header,)
    header=re.sub(re.compile(r"""\\AtBeginSection\[.*?\]\{.*?\}""",re.DOTALL),r'', header,)
    header=re.sub(re.compile(r"""\\author\{.*?\}""",re.DOTALL),r'', header,)
    header=re.sub(re.compile(r"""\\institute\{.*?\}""",re.DOTALL),r'', header,)
    header=re.sub(re.compile(r"""\\date\{.*?\}""",re.DOTALL),r'', header,)
    latexfile = open('header.tex', 'w')
    latexfile.write(header) 
    latexfile.close()

#slides = re.findall(r'\\begin{frame}(.*?)\\end{frame}', raw.read())
#items= re.findall(r"""item (.*?) """, raw.read())
slides = re.findall(r"""\\begin\{frame\}(.*?)\\end\{frame\}""",filecontent,re.DOTALL)
ideqn=0


if not os.path.exists('equations'):
    os.makedirs('equations')
for idslide,slide in enumerate(slides):
    nbeqslide=0
    newslide = prs.slides.add_slide(bullet_slide_layout)
    title= re.findall(r"""\\frametitle\s*\{(.*)\}""", slide)
    shapes = newslide.shapes
    body_shape = shapes.placeholders[1]
    
    tf = body_shape.text_frame
    tf.text = slide 
   
    #itemizes=re.findall(r"""\\begin\{itemize\}\s*(.*)\\end\{itemize\}""",slide,re.DOTALL)
    #for itemize in itemizes:
        #items=re.split(r"""\\item""",itemize,re.DOTALL)[1:]
        #for item in items:           
            #p = tf.add_paragraph()
            #p.text = 'item'
            #p.level = 1            
    
   
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    if len(title)>0:
        title_shape.text = title[0]  
        print 'add slide %s (%d over %d)'%(title[0] ,idslide,len(slides))
    equations=re.findall(r"""\\begin\{equation\}\s*(.*?)\\end\{equation\}""",slide,re.DOTALL)
    for equation in equations:
        ideqn+=1
        nbeqslide+=1
        file=os.path.join('equations','eqn%d.png'%ideqn)
        formula_as_file( equation, file, negate=False,header=header )
        left = Inches(1)
        top = Inches(2+nbeqslide)
        height =  Inches(0.5)
        pic = newslide.shapes.add_picture(file, left, top,height=height) 
        
    equations=re.findall(r"""\$\$(.*?)\$\$""",slide,re.DOTALL)
    for equation in equations:
        ideqn+=1
        nbeqslide+=1
        file=os.path.join('equations','eqn%d.png'%ideqn)
        formula_as_file( equation, file, negate=False,header=header )   
        left =  Inches(1)
        top = Inches(2+nbeqslide)
        height =  Inches(0.5)   
        pic = newslide.shapes.add_picture(file, left, top,height=height)  
        
    figures= re.findall(r"""[^%].*\\includegraphics\[.*\]\{([^\}]*)\}""", slide)
    for figure in figures:        
        left =  Inches(1)
        top = Inches(2)
        if figure[-3:]=='pdf':
            figure2=figure
            figure2=figure[:-3]+'_converted.png'
            os.system( 'convert -density 500  %s -quality 97  %s'%(figure,figure2)) 
            figure=figure2
        height = Inches(4)
        pic = newslide.shapes.add_picture(figure, left, top,height=height)        

print 'DONE'
prs.save(outputfile)
